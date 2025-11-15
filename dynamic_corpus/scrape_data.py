import os
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin
from datetime import datetime
from pathlib import Path
import fitz
import docx
from pptx import Presentation
import asyncio
from utils.extractor.file_extractor import FileExtractor

ROOT = Path(__file__).resolve().parent
DOWNLOAD_DIR = ROOT / "data/downloaded"
EXTRACT_DIR = ROOT / "data/extracted"

DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
EXTRACT_DIR.mkdir(parents=True, exist_ok=True)

TARGET_URLS = [
    # "https://legislative.gov.in/constitution-of-india",
    # "https://nhrc.nic.in",
    # "https://prsindia.org/billtrack",
    # "https://ncpcr.gov.in",
    # "https://ncw.nic.in",
    "https://prsindia.org/bills/states"
]

VALID_TYPES = ["pdf", "docx", "pptx", "txt"]

def download_file(url):
    filename = url.split("/")[-1]
    save_path = DOWNLOAD_DIR / filename

    if save_path.exists():
        print(f"Already downloaded: {filename}")
        return save_path

    print(f"Downloading: {url}")
    r = requests.get(url, timeout=20)
    r.raise_for_status()

    with open(save_path, "wb") as f:
        f.write(r.content)

    return save_path

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_docx(path):
    d = docx.Document(path)
    return "\n".join([p.text for p in d.paragraphs])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

async def extract_text(path):
    ext = path.suffix.lower()
    file_ext = FileExtractor()
    uploaded_content_data = await file_ext.extract_text(path)
    print(f"Dynamic corpus update: {uploaded_content_data.content}")
    # if ext == ".pdf":
    #     return extract_text_from_pdf(path)
    # elif ext == ".docx":
    #     return extract_text_from_docx(path)
    # elif ext == ".pptx":
    #     return extract_text_from_pptx(path)
    # elif ext == ".txt":
    #     return path.read_text()
    return uploaded_content_data.content

def save_extracted(path, text):
    out_file = EXTRACT_DIR / (path.stem + ".txt")
    out_file.write_text(text, encoding="utf-8")
    print(f"Extracted → {out_file}")

async def crawl():
    for base_url in TARGET_URLS:
        print(f"\n🔍 Crawling: {base_url}")
        r = requests.get(base_url)
        soup = BeautifulSoup(r.text, "html.parser")

        for link in soup.find_all("a", href=True):
            href = link["href"]

            # Ensure full URL
            full_url = urljoin(base_url, href)

            # Check valid types
            if any(full_url.lower().endswith(ext) for ext in VALID_TYPES):
                file_path = download_file(full_url)
                text = await extract_text(file_path)
                if text:
                    save_extracted(file_path, text)

if __name__ == "__main__":
    print("📘 Starting scrape:", datetime.now())
    asyncio.run(crawl())
    print("✔ Done.")
