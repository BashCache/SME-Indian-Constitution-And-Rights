import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from utils.extractor.base_extractor import BaseExtractor
from utils.extractor.extraction_result import ExtractionResult
from utils.extractor.gemini_summarizer import summarize_visual_content
import asyncio

class PPTXExtractor(BaseExtractor):
    """Extractor for PPTX files — includes Gemini-powered image/table summaries."""

    async def extract(self, file_path: str) -> ExtractionResult:
        slides_data, all_text = [], []

        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)

            for idx, slide in enumerate(prs.slides, start=1):
                slide_section = [f"Slide {idx}:"]
                text_blocks, table_summaries, image_summaries = [], [], []

                for shape in slide.shapes:
                    # --- 1️⃣ Text ---
                    if hasattr(shape, "text") and shape.text.strip():
                        text_blocks.append(shape.text.strip())

                    # --- 2️⃣ Tables ---
                    if shape.has_table:
                        table = shape.table
                        table_md = self._table_to_markdown(table)
                        # Optionally use Gemini for better summary
                        summary = await asyncio.run(summarize_visual_content(
                            image_path=self._table_to_image(shape, idx),
                            prompt=f"Summarize the table from slide {idx}."
                        ))
                        table_summaries.append(f"Table Data:\n{table_md}\nGemini Summary:\n{summary}")

                    # --- 3️⃣ Images ---
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_path = self._extract_image(shape, idx)
                        summary = await asyncio.run(summarize_visual_content(image_path,
                                    prompt=f"Describe image from slide {idx}."))
                        image_summaries.append(f"Image Summary:\n{summary}")

                combined_slide_text = "\n".join([
                    f"Text:\n{'\n'.join(text_blocks)}" if text_blocks else "",
                    "\n".join(table_summaries) if table_summaries else "",
                    "\n".join(image_summaries) if image_summaries else ""
                ]).strip()

                slide_text = f"Slide {idx}:\n{combined_slide_text}\n"
                slides_data.append({"slide": idx, "summary": slide_text})
                all_text.append(slide_text)

            full_text = "\n".join(all_text)
            metadata = {
                "file_type": "pptx",
                "slides": total_slides,
            }

            return ExtractionResult(content=full_text, metadata=metadata)

        except Exception as e:
            return ExtractionResult(content="", metadata={"error": str(e), "file_type": "pptx"})

    # ----------------------------------------------------------------
    def _table_to_markdown(self, table):
        rows = []
        for r in table.rows:
            cells = [c.text.strip() for c in r.cells]
            rows.append("| " + " | ".join(cells) + " |")
        return "\n".join(rows)

    def _extract_image(self, shape, slide_idx):
        image = shape.image
        image_bytes = image.blob
        tmp_dir = tempfile.gettempdir()
        img_path = os.path.join(tmp_dir, f"slide_{slide_idx}_img.png")
        with open(img_path, "wb") as f:
            f.write(image_bytes)
        return img_path

    def _table_to_image(self, shape, slide_idx):
        """(Optional) Convert table shape to image (screenshot style)."""
        return self._extract_image_placeholder(slide_idx)

    def _extract_image_placeholder(self, slide_idx):
        tmp_dir = tempfile.gettempdir()
        path = os.path.join(tmp_dir, f"slide_{slide_idx}_table.png")
        img = Image.new("RGB", (300, 100), color="white")
        img.save(path)
        return path
