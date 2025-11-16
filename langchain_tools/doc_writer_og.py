# utils/document_writer.py

import os
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document

class DocumentWriter:
    """
    Handles writing text content into various document formats (PDF, PPTX, DOCX, TXT).
    """

    @staticmethod
    def write(content: str, file_type: str, output_dir: str = "generated_docs", base_name: str = "output"):
        os.makedirs(output_dir, exist_ok=True)
        file_type = file_type.lower().strip(".")

        if file_type == "pdf":
            path = os.path.join(output_dir, f"{base_name}.pdf")
            DocumentWriter._write_pdf(content, path)
        elif file_type == "pptx":
            path = os.path.join(output_dir, f"{base_name}.pptx")
            DocumentWriter._write_pptx(content, path)
        elif file_type == "docx":
            path = os.path.join(output_dir, f"{base_name}.docx")
            DocumentWriter._write_docx(content, path)
        elif file_type == "txt":
            path = os.path.join(output_dir, f"{base_name}.txt")
            DocumentWriter._write_txt(content, path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        return path

    # ------------------------
    @staticmethod
    def _write_pdf(content: str, path: str):
        doc = SimpleDocTemplate(path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        for paragraph in content.split("\n\n"):
            story.append(Paragraph(paragraph.strip(), styles["Normal"]))
            story.append(Spacer(1, 12))

        doc.build(story)

    # ------------------------
    @staticmethod
    def _write_pptx(content: str, path: str):
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # title + content

        slides = content.split("\n\nSlide ")
        for i, slide_text in enumerate(slides, start=1):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = f"Slide {i}"
            body.text = slide_text.strip()[:1000]  # keep text reasonable per slide

        prs.save(path)

    # ------------------------
    @staticmethod
    def _write_docx(content: str, path: str):
        doc = Document()
        for para in content.split("\n\n"):
            doc.add_paragraph(para.strip())
        doc.save(path)

    # ------------------------
    @staticmethod
    def _write_txt(content: str, path: str):
        with open(path, "w", encoding="utf-8") as f:
            f.write(content.strip() + "\n")
