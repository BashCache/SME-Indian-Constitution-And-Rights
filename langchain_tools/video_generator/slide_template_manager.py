"""
Slide Template Manager
Handles PowerPoint template creation and slide generation for videos
"""

import os
from typing import List, Dict, Any, Tuple
from pathlib import Path
import tempfile
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. Slide generation will be limited.")

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: PIL/Pillow not available. Slide rendering will be limited.")

class SlideTemplateManager:
    def __init__(self):
        self.templates_dir = Path(__file__).parent / "templates"
        self.templates_dir.mkdir(exist_ok=True)
        
        # Constitutional color scheme
        self.colors = {
            'saffron': RGBColor(255, 153, 51) if PPTX_AVAILABLE else (255, 153, 51),
            'white': RGBColor(255, 255, 255) if PPTX_AVAILABLE else (255, 255, 255),
            'green': RGBColor(19, 136, 8) if PPTX_AVAILABLE else (19, 136, 8),
            'navy': RGBColor(0, 51, 102) if PPTX_AVAILABLE else (0, 51, 102),
            'gold': RGBColor(255, 215, 0) if PPTX_AVAILABLE else (255, 215, 0)
        }
    
    def create_constitutional_template(self) -> str:
        """
        Create a PowerPoint template with constitutional theme
        
        Returns:
            Path to created template file
        """
        if not PPTX_AVAILABLE:
            return self._create_basic_template()
        
        template_path = self.templates_dir / "constitutional_template.pptx"
        
        # Create new presentation
        prs = Presentation()
        
        # Slide 1: Title slide layout
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        title = slide.shapes.title
        title.text = "Constitutional Topic"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['navy']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Indian Constitution & Rights"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['saffron']
        
        # Slide 2: Content slide layout
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Main Content"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['navy']
        
        content = slide.placeholders[1]
        content.text = "• Key Point 1\n• Key Point 2\n• Key Point 3"
        
        # Save template
        prs.save(str(template_path))
        print(f"✅ Created constitutional template: {template_path}")
        
        return str(template_path)
    
    def _create_basic_template(self) -> str:
        """
        Fallback template creation without python-pptx
        """
        template_path = self.templates_dir / "basic_template.txt"
        
        template_content = """
CONSTITUTIONAL VIDEO TEMPLATE
================================

Slide 1: Title
- Title: {{title}}
- Subtitle: Indian Constitution & Rights

Slide 2-N: Content
- Header: {{section_title}}
- Content: {{content}}

Color Scheme:
- Primary: Navy Blue (0, 51, 102)
- Secondary: Saffron (255, 153, 51)
- Accent: Green (19, 136, 8)
"""
        
        with open(template_path, 'w') as f:
            f.write(template_content)
        
        return str(template_path)
    
    def generate_slides_from_script(self, script_data: Dict[str, Any], output_dir: str) -> List[str]:
        """
        Generate individual slide images from script data
        
        Args:
            script_data: Dictionary containing slide content
            output_dir: Directory to save slide images
        
        Returns:
            List of slide image paths
        """
        output_path = Path(output_dir)
        output_path.mkdir(exist_ok=True)
        
        slide_images = []
        
        # Generate slides based on script segments
        segments = script_data.get('segments', [])
        
        for i, segment in enumerate(segments):
            slide_path = output_path / f"slide_{i+1}_{segment.get('slide_id', 'content')}.png"
            
            if PIL_AVAILABLE:
                success = self._create_slide_image(
                    title=segment.get('title', f"Slide {i+1}"),
                    content=segment.get('content', ''),
                    slide_type=segment.get('type', 'content'),
                    output_path=str(slide_path)
                )
                
                if success:
                    slide_images.append(str(slide_path))
            else:
                # Create placeholder image info
                slide_images.append(str(slide_path))
                print(f"Generated placeholder for: {slide_path}")
        
        return slide_images
    
    def _create_slide_image(self, title: str, content: str, slide_type: str, output_path: str) -> bool:
        """
        Create individual slide image using PIL
        
        Args:
            title: Slide title
            content: Slide content
            slide_type: Type of slide (title, content, conclusion)
            output_path: Where to save the image
        
        Returns:
            Success status
        """
        try:
            # Create image with 16:9 aspect ratio
            width, height = 1920, 1080
            image = Image.new('RGB', (width, height), color=self.colors['white'])
            draw = ImageDraw.Draw(image)
            
            # Try to load fonts (fallback to default if not available)
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 72)
                content_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 48)
                subtitle_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 36)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                subtitle_font = ImageFont.load_default()
            
            if slide_type == 'title':
                # Title slide layout
                # Draw title
                title_bbox = draw.textbbox((0, 0), title, font=title_font)
                title_x = (width - (title_bbox[2] - title_bbox[0])) // 2
                title_y = height // 3
                draw.text((title_x, title_y), title, font=title_font, fill=self.colors['navy'])
                
                # Draw subtitle
                subtitle = "Indian Constitution & Rights"
                subtitle_bbox = draw.textbbox((0, 0), subtitle, font=subtitle_font)
                subtitle_x = (width - (subtitle_bbox[2] - subtitle_bbox[0])) // 2
                subtitle_y = title_y + 120
                draw.text((subtitle_x, subtitle_y), subtitle, font=subtitle_font, fill=self.colors['saffron'])
                
            else:
                # Content slide layout
                # Draw title
                title_bbox = draw.textbbox((0, 0), title, font=title_font)
                title_x = 100
                title_y = 100
                draw.text((title_x, title_y), title, font=title_font, fill=self.colors['navy'])
                
                # Draw content with line wrapping
                content_lines = self._wrap_text(content, content_font, width - 200)
                content_y = title_y + 150
                
                for line in content_lines:
                    draw.text((title_x, content_y), line, font=content_font, fill=(50, 50, 50))
                    content_y += 60
            
            # Add decorative elements
            # Top border in saffron
            draw.rectangle([(0, 0), (width, 20)], fill=self.colors['saffron'])
            # Bottom border in green
            draw.rectangle([(0, height-20), (width, height)], fill=self.colors['green'])
            
            # Save image
            image.save(output_path, 'PNG')
            print(f"✅ Generated slide: {output_path}")
            return True
            
        except Exception as e:
            print(f"❌ Error creating slide image: {str(e)}")
            return False
    
    def _wrap_text(self, text: str, font, max_width: int) -> List[str]:
        """
        Wrap text to fit within specified width
        """
        if not PIL_AVAILABLE:
            return text.split('\n')
        
        lines = []
        words = text.split()
        current_line = ""
        
        for word in words:
            test_line = current_line + (" " if current_line else "") + word
            bbox = ImageDraw.Draw(Image.new('RGB', (1, 1))).textbbox((0, 0), test_line, font=font)
            line_width = bbox[2] - bbox[0]
            
            if line_width <= max_width:
                current_line = test_line
            else:
                if current_line:
                    lines.append(current_line)
                current_line = word
        
        if current_line:
            lines.append(current_line)
        
        return lines