import os
import json
import io
from pptx import Presentation
from PIL import Image
import pytesseract
import pandas as pd
import subprocess
import shutil
import tempfile

class PPTExtractor:
    def __init__(self, ppt_folder, output_folder):
        """
        Initialize the PPT extractor
        
        Args:
            ppt_folder (str): Path to folder containing PPT files
            output_folder (str): Path to store extracted text files
        """
        self.ppt_folder = ppt_folder
        self.output_folder = output_folder
        os.makedirs(self.output_folder, exist_ok=True)
        
    def extract_table_data(self, table):
        """Extract data from a table and convert to text format"""
        rows = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                text = cell.text.strip()
                row_data.append(text)
            rows.append(row_data)
        
        # Convert to pandas DataFrame for easier handling
        df = pd.DataFrame(rows[1:], columns=rows[0] if rows else None)
        return df.to_string()

    def preprocess_image(self, image):
        """Preprocess image for better OCR results"""
        try:
            # Convert to RGB if not already
            if image.mode != 'RGB':
                image = image.convert('RGB')

            # Convert to grayscale
            image = image.convert('L')
            
            from PIL import ImageEnhance, ImageOps
            
            # Apply binary thresholding
            threshold = 200
            image = image.point(lambda x: 0 if x < threshold else 255, '1')
            
            # Increase contrast
            image = ImageOps.autocontrast(image)
            
            # Denoise
            image = image.filter(ImageEnhance.SMOOTH)
            
            # Resize if too small (minimum 2000px on longest side for better OCR)
            if image.width < 2000 or image.height < 2000:
                ratio = max(2000/image.width, 2000/image.height)
                new_size = (int(image.width * ratio), int(image.height * ratio))
                image = image.resize(new_size, Image.Resampling.LANCZOS)
            
            return image
        except Exception as e:
            # print(f"Error in image preprocessing: {e}")
            return image

    def extract_image_text(self, image):
        """Extract text from image using OCR"""
        try:
            # Set Tesseract path
            pytesseract.pytesseract.tesseract_cmd = r'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'
            # Working on Windows, remove above forother OS
            
            # Check if tesseract is installed and accessible
            try:
                pytesseract.get_tesseract_version()
            except Exception as e:
                print(f"Tesseract not properly configured: {e}")
                print("Please ensure Tesseract is installed and the path is correct")
                return ""

            # Preprocess the image
            processed_image = self.preprocess_image(image)
            
            # Save debug image to check preprocessing
            debug_dir = os.path.join(self.output_folder, "debug_images")
            os.makedirs(debug_dir, exist_ok=True)
            debug_path = os.path.join(debug_dir, f"debug_image_{len(os.listdir(debug_dir))}.png")
            processed_image.save(debug_path)
            print(f"Saved debug image to: {debug_path}")
            
            # Try OCR with different configurations
            configs = [
                '--oem 3 --psm 6 -l eng --dpi 300',  # Default
                '--oem 3 --psm 1 -l eng --dpi 300',  # Automatic page segmentation
                '--oem 3 --psm 3 -l eng --dpi 300',  # Fully automatic page segmentation
                '--oem 3 --psm 4 -l eng --dpi 300',  # Assume single column of text
                '--oem 3 --psm 11 -l eng --dpi 300', # Sparse text - no orientation assumption
                '--oem 1 --psm 6 -l eng --dpi 300',  # Legacy engine with default PSM
            ]
            
            for config in configs:
                text = pytesseract.image_to_string(processed_image, config=config)
                if text.strip():
                    print(f"Successfully extracted text with config: {config}")
                    return text.strip()
            
            print("No text could be extracted from the image")
            return ""
            
        except Exception as e:
            print(f"Error extracting text from image: {str(e)}")
            return ""

    def process_slide(self, slide):
        """Process a single slide and extract all content"""
        content = {
            "text": [],
            "tables": [],
            "image_text": []
        }
        
        # Extract text from shapes
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text and len(text) > 20:
                    content["text"].append(text)
            
            # Extract tables
            if shape.has_table:
                table_text = self.extract_table_data(shape.table)
                if table_text and len(table_text.strip()) > 100:
                    content["tables"].append(table_text)
                    # Add table text to main text as well
                    content["text"].append(f"Table content:\n{table_text}")
            
            # Extract text from images using OCR
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                if image:
                    # Process the image directly from memory
                    img = Image.open(io.BytesIO(image.blob))
                    image_text = self.extract_image_text(img)
                    img.close()
                    
                    if image_text and len(image_text.strip()) > 100:
                        content["image_text"].append(image_text)
                        # Add image text to main text as well
                        content["text"].append(f"Image text:\n{image_text}")
        
        return content

    def process_presentation(self, ppt_path):
        """Process entire presentation and save extracted content"""
        try:
            presentation = Presentation(ppt_path)
            all_content = []
            
            for i, slide in enumerate(presentation.slides, 1):
                slide_content = self.process_slide(slide)
                slide_content["slide_number"] = i
                all_content.append(slide_content)
            
            # Save extracted content
            base_name = os.path.splitext(os.path.basename(ppt_path))[0]
            output_filename = os.path.join(self.output_folder, f"{base_name}_extracted.json")
            with open(output_filename, 'w', encoding='utf-8') as f:
                json.dump(all_content, f, indent=2, ensure_ascii=False)
            
            print(f"Successfully processed: {os.path.basename(ppt_path)}")
            
        except Exception as e:
            print(f"Error processing {os.path.basename(ppt_path)}: {e}")

    def process_all_presentations(self):
        """Process all PPT files in the specified folder"""
        ppt_files = [f for f in os.listdir(self.ppt_folder) 
                    if f.endswith(('.ppt', '.pptx')) and 
                    os.path.isfile(os.path.join(self.ppt_folder, f))]
        
        if not ppt_files:
            print(f"No PPT files found in {self.ppt_folder}")
            return
        
        for ppt_file in ppt_files:
            full_path = os.path.join(self.ppt_folder, ppt_file)
            self.process_presentation(full_path)

if __name__ == "__main__":
    # Set up paths
    ppt_folder = os.path.join("data", "ppt")
    output_folder = os.path.join("extracted_data", "ppt_extracted")

    # Initialize and run extractor
    extractor = PPTExtractor(ppt_folder, output_folder)
    extractor.process_all_presentations()