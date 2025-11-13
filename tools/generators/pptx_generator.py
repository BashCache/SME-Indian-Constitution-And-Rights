"""
PPTX Generator: Generates Microsoft PowerPoint presentations from text content.
"""

from typing import Optional
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

from .document_generator import DocumentGenerator


class PptxGenerator(DocumentGenerator):
    """Generator for PPTX (Microsoft PowerPoint) presentations"""
    
    def generate(self, content: str, output_path: str, title: Optional[str] = None,
                author: Optional[str] = None, subject: Optional[str] = None) -> str:
        """Generate a PPTX presentation from text content"""
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Set presentation properties
        if title:
            prs.core_properties.title = title
        if author:
            prs.core_properties.author = author
        if subject:
            prs.core_properties.subject = subject
        
        # Split content into slides
        slides_content = self.split_into_slides(content, max_slide_length=300)
        
        # Add title slide if title is provided
        if title:
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            if subject:
                subtitle_shape.text = subject
            elif author:
                subtitle_shape.text = f"By {author}"
        
        # Add content slides
        for slide_content in slides_content:
            # Use blank layout for content slides
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add text box
            left = PptxInches(0.5)
            top = PptxInches(1)
            width = PptxInches(9)
            height = PptxInches(5.5)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            # Split slide content into paragraphs
            paragraphs = self.split_into_paragraphs(slide_content)
            
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = text_frame.paragraphs[0]
                    p.text = para_text
                    p.font.size = PptxPt(18)
                    p.font.bold = True
                else:
                    p = text_frame.add_paragraph()
                    p.text = para_text
                    p.font.size = PptxPt(14)
                    p.level = 0
                
                p.alignment = PP_ALIGN.LEFT
                p.space_after = PptxPt(12)
        
        # Save presentation
        prs.save(output_path)
        return output_path

